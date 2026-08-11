from __future__ import annotations
from io import BytesIO
from pathlib import Path
from uuid import UUID
from datetime import datetime
import json, os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from app.repositories.finance.gl_transaction_repository import GLTransactionRepository
from app.services.finance.reporting_service import ReportingService
from app.services.finance.analytics_service import AnalyticsService

ARTIFACT_ROOT=Path('generated_artifacts')

class BoardPackService:
    def __init__(self, session: AsyncSession):
        self.session=session; self.reporting=ReportingService(GLTransactionRepository(session)); self.analytics=AnalyticsService(session)

    async def _data(self,company_id,forecast_run_id):
        pnl=await self.reporting.pnl(company_id); bs=await self.reporting.balance_sheet(company_id); monthly=await self.reporting.monthly_actuals(company_id); branches=await self.reporting.branch_comparison(company_id); analytics=await self.analytics.overview(company_id)
        forecast=None
        if forecast_run_id:
            forecast=(await self.session.execute(text('SELECT summary,result_payload FROM public.forecast_model_runs WHERE company_id=:c AND id=:i'),{'c':company_id,'i':forecast_run_id})).mappings().first()
        return {'pnl':pnl,'bs':bs,'monthly':monthly,'branches':branches,'analytics':analytics,'forecast':dict(forecast) if forecast else None}

    def _chart(self,monthly,path):
        if not monthly:return None
        periods=[str(x['month'])[:7] for x in monthly]; revenue=[float(x['revenue']) for x in monthly]; profit=[float(x['net_profit']) for x in monthly]
        fig,ax=plt.subplots(figsize=(9,4.2)); ax.plot(periods,revenue,marker='o',label='Revenue'); ax.plot(periods,profit,marker='o',label='Net Profit'); ax.legend(); ax.grid(alpha=.25); ax.tick_params(axis='x',rotation=45); fig.tight_layout(); fig.savefig(path,dpi=160); plt.close(fig); return path

    async def generate(self,company,request):
        data=await self._data(company.id,request.forecast_run_id)
        run_id=(await self.session.execute(text('''INSERT INTO public.board_pack_runs(company_id,pack_name,reporting_period,forecast_run_id,selected_sections,commentary)
          VALUES (:c,:n,:p,:f,CAST(:s AS jsonb),CAST(:m AS jsonb)) RETURNING id'''),{'c':company.id,'n':request.pack_name,'p':request.reporting_period,'f':request.forecast_run_id,'s':json.dumps(request.sections),'m':json.dumps(request.model_dump())})).scalar_one()
        folder=ARTIFACT_ROOT/str(company.id)/str(run_id); folder.mkdir(parents=True,exist_ok=True); chart=self._chart(data['monthly'],folder/'monthly.png')
        artifacts=[]
        for fmt in request.formats:
            if fmt=='pptx': path=self._pptx(folder,company,request,data,chart)
            elif fmt=='pdf': path=self._pdf(folder,company,request,data,chart)
            elif fmt=='xlsx': path=self._xlsx(folder,company,request,data)
            else: continue
            aid=(await self.session.execute(text('''INSERT INTO public.generated_artifacts(company_id,board_pack_run_id,artifact_type,file_name,storage_path,file_size_bytes)
              VALUES (:c,:r,:t,:n,:p,:z) RETURNING id'''),{'c':company.id,'r':run_id,'t':fmt,'n':path.name,'p':str(path.resolve()),'z':path.stat().st_size})).scalar_one()
            artifacts.append({'id':aid,'artifact_type':fmt,'file_name':path.name,'download_url':f'/api/v1/board-packs/artifacts/{aid}/download','file_size_bytes':path.stat().st_size})
        await self.session.commit(); return artifacts

    def _pptx(self,folder,company,request,data,chart):
        prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
        s=prs.slides.add_slide(prs.slide_layouts[0]); s.shapes.title.text=request.pack_name; s.placeholders[1].text=f'{company.legal_name}\n{request.reporting_period}\nConfidential'
        def bullet(title,items):
            sl=prs.slides.add_slide(prs.slide_layouts[1]); sl.shapes.title.text=title; tf=sl.placeholders[1].text_frame; tf.clear()
            for i,item in enumerate(items): p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=str(item); p.font.size=Pt(20)
        pnl=data['pnl']; bs=data['bs']
        bullet('Executive Summary',[f"Revenue: {pnl.revenue:,.0f}",f"Net profit: {pnl.net_profit:,.0f}",f"Total assets: {bs.total_assets:,.0f}",f"Closing equity: {bs.equity:,.0f}",request.management_outlook or 'Management outlook not supplied.'])
        if chart:
            sl=prs.slides.add_slide(prs.slide_layouts[5]); sl.shapes.title.text='Monthly Performance'; sl.shapes.add_picture(str(chart),Inches(1),Inches(1.4),width=Inches(11.2))
        bullet('Strategic Priorities',[request.strategic_priorities or 'Not supplied',request.principal_risks or 'Risks not supplied',request.decisions_required or 'No decisions recorded'])
        if data['forecast']: bullet('Forecast & Three-Way Outlook',[f"Forecast revenue: {data['forecast']['summary'].get('forecast_revenue',0):,.0f}",f"Forecast EBITDA: {data['forecast']['summary'].get('forecast_ebitda',0):,.0f}",f"Closing cash: {data['forecast']['summary'].get('closing_cash',0):,.0f}",f"Balanced: {data['forecast']['summary'].get('balanced')}"])
        path=folder/f"{request.pack_name.replace(' ','_')}.pptx"; prs.save(path); return path

    def _pdf(self,folder,company,request,data,chart):
        path=folder/f"{request.pack_name.replace(' ','_')}.pdf"; styles=getSampleStyleSheet(); story=[Paragraph(request.pack_name,styles['Title']),Paragraph(company.legal_name,styles['Heading2']),Paragraph(request.reporting_period,styles['Normal']),Spacer(1,18)]
        pnl=data['pnl']; bs=data['bs']; table=Table([['Measure','Amount'],['Revenue',f'{pnl.revenue:,.0f}'],['Gross Profit',f'{pnl.gross_profit:,.0f}'],['Net Profit',f'{pnl.net_profit:,.0f}'],['Total Assets',f'{bs.total_assets:,.0f}'],['Equity',f'{bs.equity:,.0f}']]); table.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#17345B')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),.5,colors.grey),('PADDING',(0,0),(-1,-1),8)])); story+=[table,Spacer(1,20)]
        if chart: story += [Image(str(chart),width=500,height=230),PageBreak()]
        for title,body in [('Management Outlook',request.management_outlook),('Strategic Priorities',request.strategic_priorities),('Principal Risks',request.principal_risks),('Decisions Required',request.decisions_required)]: story += [Paragraph(title,styles['Heading1']),Paragraph(body or 'Not supplied',styles['BodyText']),Spacer(1,12)]
        SimpleDocTemplate(str(path),pagesize=landscape(A4),rightMargin=30,leftMargin=30,topMargin=30,bottomMargin=30).build(story); return path

    def _xlsx(self,folder,company,request,data):
        path=folder/f"{request.pack_name.replace(' ','_')}.xlsx"; wb=Workbook(); ws=wb.active; ws.title='Executive Summary'; ws.append([request.pack_name]); ws.append([company.legal_name,request.reporting_period]); ws.append([]); ws.append(['Measure','Amount']); pnl=data['pnl']; bs=data['bs'];
        for row in [('Revenue',pnl.revenue),('Gross Profit',pnl.gross_profit),('Net Profit',pnl.net_profit),('Total Assets',bs.total_assets),('Total Liabilities',bs.total_liabilities),('Equity',bs.equity)]:ws.append(row)
        ws['A1'].font=Font(bold=True,size=18); ws['A4'].fill=PatternFill('solid',fgColor='17345B'); ws['B4'].fill=PatternFill('solid',fgColor='17345B'); ws['A4'].font=ws['B4'].font=Font(color='FFFFFF',bold=True)
        m=wb.create_sheet('Monthly Actuals'); m.append(['Month','Revenue','Gross Profit','Operating Expenses','EBIT','Net Profit']);
        for r in data['monthly']:m.append([r['month'],r['revenue'],r['gross_profit'],r['operating_expenses'],r['ebit'],r['net_profit']])
        wb.save(path); return path
