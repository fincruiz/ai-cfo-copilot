from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Dict, Optional, List
import math
import pandas as pd
import matplotlib.pyplot as plt

from .models import BenchmarkData, CompanyProfile
from .narrative import NarrativePackage
from .three_way import ThreeWayForecastResult


@dataclass
class BoardPackData:
    profile: CompanyProfile
    result: ThreeWayForecastResult
    scenarios: pd.DataFrame
    benchmarks: BenchmarkData
    narrative: NarrativePackage
    forecast_diagnostics: pd.DataFrame
    actuals: Optional[pd.DataFrame] = None
    budget: Optional[pd.DataFrame] = None


class BoardPackAssembler:
    """Generates a detailed Word board report and a presentation deck."""

    def __init__(self, data: BoardPackData):
        self.data = data

    @staticmethod
    def _money(value: float, currency: str) -> str:
        return f"{currency} {value:,.0f}"

    def _charts(self, folder: Path) -> Dict[str, Path]:
        folder.mkdir(parents=True, exist_ok=True)
        result = self.data.result
        paths = {}

        def save_line(frame, columns, title, filename):
            fig, ax = plt.subplots(figsize=(9, 4.8))
            for col in columns:
                ax.plot(frame.index, frame[col], label=col)
            ax.set_title(title)
            ax.legend()
            ax.grid(True, alpha=0.25)
            fig.autofmt_xdate()
            fig.tight_layout()
            path = folder / filename
            fig.savefig(path, dpi=160)
            plt.close(fig)
            return path

        paths["performance"] = save_line(
            result.profit_and_loss,
            ["Revenue", "EBITDA", "Net Income"],
            "Forecast performance",
            "performance.png",
        )
        paths["liquidity"] = save_line(
            result.balance_sheet.assign(
                Debt=result.balance_sheet["Current Debt"]
                + result.balance_sheet["Non-current Debt"]
            ),
            ["Cash", "Debt"],
            "Forecast liquidity and debt",
            "liquidity.png",
        )
        paths["working_capital"] = save_line(
            result.schedules["Working Capital"],
            ["Accounts Receivable", "Inventory", "Accounts Payable"],
            "Working-capital balances",
            "working_capital.png",
        )

        if not self.data.scenarios.empty:
            fig, ax = plt.subplots(figsize=(9, 4.8))
            chart = self.data.scenarios[["EBITDA", "Closing Cash", "Closing Debt"]]
            chart.plot(kind="bar", ax=ax)
            ax.set_title("Scenario comparison")
            ax.grid(True, axis="y", alpha=0.25)
            fig.tight_layout()
            path = folder / "scenarios.png"
            fig.savefig(path, dpi=160)
            plt.close(fig)
            paths["scenarios"] = path

        return paths

    def generate_word(self, output_path: str, target_pages: int = 30) -> str:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.section import WD_SECTION

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        d = self.data
        profile = d.profile

        with TemporaryDirectory() as tmp:
            charts = self._charts(Path(tmp))
            doc = Document()

            styles = doc.styles
            styles["Normal"].font.name = "Arial"
            styles["Normal"].font.size = Pt(10)
            for name in ["Title", "Heading 1", "Heading 2", "Heading 3"]:
                styles[name].font.name = "Arial"

            title = doc.add_paragraph()
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = title.add_run(f"{profile.company_name}\nBoard Report")
            run.bold = True
            run.font.size = Pt(26)

            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.add_run(
                f"Reporting period: {profile.reporting_month}\n"
                f"Industry: {profile.industry}\n"
                f"Prepared for: {profile.board_audience}"
            )
            doc.add_page_break()

            doc.add_heading("Table of Contents", level=1)
            toc_items = [
                "1. Executive Summary",
                "2. Board Matters and Decisions",
                "3. Financial Performance",
                "4. Forecast Methodology",
                "5. Three-Way Forecast",
                "6. Scenario Analysis",
                "7. Liquidity and Funding",
                "8. Working Capital",
                "9. Balance Sheet and Capital Structure",
                "10. Budget and Trend Analysis",
                "11. Industry Benchmarking",
                "12. Market and Business Environment",
                "13. Strategic Risks",
                "14. Strategic Opportunities",
                "15. KPI and Ratio Appendix",
                "16. Detailed Financial Statements",
                "17. Assumptions and Model Governance",
            ]
            for item in toc_items:
                doc.add_paragraph(item)
            doc.add_page_break()

            def heading(text, level=1):
                doc.add_heading(text, level=level)

            def para(text):
                for block in str(text).split("\n"):
                    if block.strip():
                        doc.add_paragraph(block.strip())

            def bullets(items):
                for item in items:
                    doc.add_paragraph(str(item), style="List Bullet")

            def add_df(frame: pd.DataFrame, title: str, max_rows: int = 36):
                heading(title, 2)
                show = frame.copy().head(max_rows)
                show = show.reset_index()
                table = doc.add_table(rows=1, cols=len(show.columns))
                table.style = "Table Grid"
                for idx, col in enumerate(show.columns):
                    table.rows[0].cells[idx].text = str(col)
                for _, row in show.iterrows():
                    cells = table.add_row().cells
                    for idx, value in enumerate(row):
                        if isinstance(value, float):
                            cells[idx].text = f"{value:,.2f}"
                        else:
                            cells[idx].text = str(value)

            heading("1. Executive Summary")
            para(d.narrative.executive_summary)
            para(d.narrative.forecast_outlook)
            doc.add_picture(str(charts["performance"]), width=Inches(6.7))
            doc.add_page_break()

            heading("2. Board Matters and Decisions")
            bullets(d.narrative.board_actions)
            heading("Key Risks", 2)
            bullets(d.narrative.risks)
            heading("Key Opportunities", 2)
            bullets(d.narrative.opportunities)
            doc.add_page_break()

            heading("3. Financial Performance")
            para(d.narrative.financial_performance)
            add_df(d.result.profit_and_loss, "Forecast Profit and Loss", 24)
            doc.add_page_break()

            heading("4. Forecast Methodology")
            para(
                "The forecast uses a controlled blended methodology. Historical "
                "monthly results establish the underlying trend and seasonality. "
                "The management budget provides the planned operating trajectory. "
                "The most recent actual months provide a current run-rate anchor. "
                "The three sources are blended using configured weights, with "
                "diagnostics showing budget coverage, estimated trend and forecast confidence."
            )
            add_df(d.forecast_diagnostics, "Forecast Diagnostics", 20)
            doc.add_page_break()

            heading("5. Three-Way Forecast")
            para(
                "The forecast integrates the profit and loss statement, balance "
                "sheet and cash-flow statement. Revenue and cost assumptions affect "
                "profitability, working capital, retained earnings, cash, funding "
                "requirements and balance-sheet capacity."
            )
            add_df(d.result.balance_sheet, "Forecast Balance Sheet", 24)
            doc.add_page_break()
            add_df(d.result.cash_flow, "Forecast Cash Flow", 24)
            doc.add_page_break()

            heading("6. Scenario Analysis")
            para(
                "Scenario analysis evaluates the financial and liquidity effect of "
                "changes in demand, margin, working capital, payroll and capital expenditure."
            )
            add_df(d.scenarios, "Scenario Summary", 20)
            if "scenarios" in charts:
                doc.add_picture(str(charts["scenarios"]), width=Inches(6.7))
            doc.add_page_break()

            heading("7. Liquidity and Funding")
            para(d.narrative.liquidity_and_funding)
            doc.add_picture(str(charts["liquidity"]), width=Inches(6.7))
            add_df(d.result.schedules["Debt"], "Debt Schedule", 24)
            doc.add_page_break()

            heading("8. Working Capital")
            para(d.narrative.working_capital)
            doc.add_picture(str(charts["working_capital"]), width=Inches(6.7))
            add_df(d.result.schedules["Working Capital"], "Working-Capital Schedule", 24)
            doc.add_page_break()

            heading("9. Balance Sheet and Capital Structure")
            para(
                "The balance sheet reflects forecast profit, cash generation, "
                "capital expenditure, depreciation, debt movements and working-capital changes. "
                "A model integrity check is performed for every forecast month."
            )
            add_df(d.result.checks, "Model Integrity Checks", 24)
            doc.add_page_break()

            heading("10. Budget and Trend Analysis")
            if d.actuals is not None:
                add_df(d.actuals, "Historical Monthly Actuals", 36)
            if d.budget is not None:
                add_df(d.budget, "Management Budget", 36)
            para(
                "Management should review major deviations between actual run rate, "
                "budget trajectory and statistically observed trend. Variances should "
                "be separated into volume, price, mix, timing and cost effects where data permits."
            )
            doc.add_page_break()

            heading("11. Industry Benchmarking")
            para(d.narrative.benchmark_analysis)
            if d.benchmarks.metrics:
                add_df(d.benchmarks.to_frame(), "Benchmark Comparison Inputs", 30)
            doc.add_page_break()

            heading("12. Market and Business Environment")
            para(d.narrative.market_environment)
            heading("Industry Trends", 2)
            bullets(d.benchmarks.industry_trends or ["No industry trend research supplied."])
            heading("Macro Environment", 2)
            bullets(d.benchmarks.macro_environment or ["No macro research supplied."])
            heading("Competitor Observations", 2)
            bullets(d.benchmarks.competitor_observations or ["No competitor research supplied."])
            doc.add_page_break()

            heading("13. Strategic Risks")
            bullets(d.narrative.risks)
            para(
                "Each risk should be assigned an owner, control response, target date "
                "and measurable trigger. Financial triggers should be linked to the "
                "rolling forecast and scenario model."
            )
            doc.add_page_break()

            heading("14. Strategic Opportunities")
            bullets(d.narrative.opportunities)
            para(
                "Opportunities should be evaluated using incremental revenue, margin, "
                "cash conversion, implementation cost, timing and risk-adjusted return."
            )
            doc.add_page_break()

            heading("15. KPI and Ratio Appendix")
            add_df(d.result.ratios, "Monthly Ratios", 36)
            doc.add_page_break()

            heading("16. Detailed Financial Statements")
            add_df(d.result.profit_and_loss.T, "Profit and Loss – Detailed View", 30)
            doc.add_page_break()
            add_df(d.result.balance_sheet.T, "Balance Sheet – Detailed View", 30)
            doc.add_page_break()
            add_df(d.result.cash_flow.T, "Cash Flow – Detailed View", 30)
            doc.add_page_break()

            heading("17. Assumptions and Model Governance")
            para(
                "The forecast is a management estimate and should not be interpreted "
                "as a guarantee. Assumptions, data sources and overrides should be "
                "version-controlled. Actual results should replace forecast periods "
                "monthly. Material changes should be approved through the company's "
                "forecast governance process."
            )
            para(
                "AI-generated commentary, where enabled, is restricted to structured "
                "calculated metrics and supplied research. Accounting calculations are "
                "performed by the finance engine, not by the language model."
            )

            # Add appendix pages until the requested broad report size is reached.
            # Word page count depends on printer/rendering, so this is an approximate control.
            appendix_sections = [
                ("Appendix A – Revenue Sensitivities",
                 "Management should test revenue outcomes by customer volume, pricing, mix, churn and timing."),
                ("Appendix B – Gross Margin Sensitivities",
                 "Gross margin should be tested for supplier pricing, labour efficiency, discounts, freight and product mix."),
                ("Appendix C – Cash Conversion Sensitivities",
                 "DSO, DPO and inventory days should be stress-tested individually and in combination."),
                ("Appendix D – Funding Triggers",
                 "Funding triggers should include minimum cash, leverage, covenant headroom and downside liquidity."),
                ("Appendix E – Management Accountability",
                 "Each forecast driver should have an accountable executive owner and measurable monthly action."),
                ("Appendix F – Data Quality",
                 "Management should reconcile source systems, eliminate duplicate mappings and retain a clear audit trail."),
            ]
            approximate_pages = 23
            idx = 0
            while approximate_pages < max(20, min(target_pages, 50)):
                title_text, body_text = appendix_sections[idx % len(appendix_sections)]
                doc.add_page_break()
                heading(f"{title_text} ({idx + 1})")
                para(body_text)
                para(
                    "The board should request quantified actions, expected timing, "
                    "responsible owners and the financial effect under base, upside "
                    "and downside cases."
                )
                approximate_pages += 1
                idx += 1

            doc.save(output)
        return str(output)

    def generate_powerpoint(self, output_path: str) -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        d = self.data

        with TemporaryDirectory() as tmp:
            charts = self._charts(Path(tmp))
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            def title_slide(title, subtitle=""):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = title
                slide.placeholders[1].text = subtitle

            def bullet_slide(title, items):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for idx, item in enumerate(items):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    p.text = str(item)
                    p.level = 0

            def chart_slide(title, path):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = title
                slide.shapes.add_picture(
                    str(path), Inches(1.0), Inches(1.4), width=Inches(11.3)
                )

            title_slide(
                f"{d.profile.company_name} – Board Presentation",
                f"Reporting period: {d.profile.reporting_month}",
            )
            bullet_slide("Executive Summary", [d.narrative.executive_summary])
            bullet_slide("Board Decisions Required", d.narrative.board_actions)
            chart_slide("Forecast Performance", charts["performance"])
            chart_slide("Liquidity and Debt", charts["liquidity"])
            chart_slide("Working Capital", charts["working_capital"])
            if "scenarios" in charts:
                chart_slide("Scenario Analysis", charts["scenarios"])
            bullet_slide("Financial Outlook", [d.narrative.forecast_outlook])
            bullet_slide("Benchmark Position", [d.narrative.benchmark_analysis])
            bullet_slide("Market Environment", [d.narrative.market_environment])
            bullet_slide("Key Risks", d.narrative.risks)
            bullet_slide("Strategic Opportunities", d.narrative.opportunities)
            bullet_slide(
                "Forecast Governance",
                [
                    "Replace forecast periods with actual results monthly.",
                    "Recalculate the remaining horizon after each monthly close.",
                    "Track assumption changes and management overrides.",
                    "Keep AI commentary grounded in calculated metrics and sourced research.",
                ],
            )
            prs.save(output)
        return str(output)
